from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "31 Maart"
subtitle.text = "Gemaakt door: Yasha en Olivier"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Wat gebeurde er met de nsb’ers in Ede"
bullet_points = [
    "NSB'ers sympathiseerden met de Duitse bezetter",
    "Na de oorlog werden zij gezien als verraders",
    "In Ede werden zij opgepakt en vastgehouden in een kamp",
    "De straffen waren bedoeld als vergelding",
    "Sommigen werden publiekelijk vernederd"
]
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.level = 0

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Hoe en wanneer werd Ede bevrijd tijdens de WO2?"
bullet_points = [
    "Ede werd bevrijd op 17 april 1945",
    "Polar Bear divisie en Britse tanks rukten op",
    "300 Duitse militairen waren op dat moment aanwezig",
    "Het Glosters Bataljon vertrok richting Renkum",
    "Het Essex Bataljon stak de Rijn over met landingsvaartuigen",
    "De A-compagnie van het Duke of Wellington Bataljon zuiverde Wageningen-Zuid en -West",
    "De bevrijding van Ede was een keerpunt in de bevrijding van Nederland"
]
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.level = 0

prs.save("31_maart_presentatie.pptx")
